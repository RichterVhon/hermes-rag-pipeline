import re
import requests
import psycopg2
import os
import sys
import concurrent.futures

# TODO (future): let an AI classify the content type (FAQ / multi-column table /
# legal document / transcript / etc.) and pick a chunking strategy accordingly,
# instead of always using the flat ~800-char paragraph splitter below.

def chunk_text(text, max_chars=800):
    paragraphs = [p.strip() for p in text.split(chr(10)) if p.strip() and len(p.strip()) > 40]
    chunks, current = [], ''
    for para in paragraphs:
        if len(current) + len(para) <= max_chars:
            current += para + ' '
        else:
            if current:
                chunks.append(current.strip())
            current = para + ' '
    if current:
        chunks.append(current.strip())
    return chunks

def chunk_faq(text, max_chars=800):
    lines = [l.strip() for l in text.split(chr(10)) if l.strip()]
    chunks, current = [], []
    for line in lines:
        is_question_start = line.endswith("?") or re.match(r"^(Q:|Question:)", line, re.IGNORECASE)
        if is_question_start and current:
            chunks.append(chr(10).join(current))
            current = [line]
        else:
            current.append(line)
    if current:
        chunks.append(chr(10).join(current))
    return chunks

def chunk_legal(text, max_chars=1200):
    pattern = r"(?=^\s*(?:Section\s+\d+|Article\s+\d+|\d+(?:\.\d+)*\.?\s))"
    parts = re.split(pattern, text, flags=re.MULTILINE)
    return [p.strip() for p in parts if p.strip()]

def chunk_transcript(text, max_chars=1000):
    pattern = r"(?=^[A-Z][A-Za-z .]{0,30}:\s)"
    turns = [t.strip() for t in re.split(pattern, text, flags=re.MULTILINE) if t.strip()]
    chunks, current = [], ""
    for turn in turns:
        if len(current) + len(turn) <= max_chars:
            current += turn + chr(10)
        else:
            if current:
                chunks.append(current.strip())
            current = turn + chr(10)
    if current:
        chunks.append(current.strip())
    return chunks

def chunk_procedural(text, max_chars=800):
    pattern = r"(?=^\s*(?:Step\s+\d+|\d+[.)]\s))"
    parts = re.split(pattern, text, flags=re.MULTILINE)
    return [p.strip() for p in parts if p.strip()]

def chunk_sectioned(text, max_chars=1200):
    lines = [l.strip() for l in text.split(chr(10)) if l.strip()]
    chunks, current = [], []
    for line in lines:
        looks_like_header = len(line) < 60 and not line.endswith(".") and not line.endswith(",")
        if looks_like_header and current:
            chunks.append(chr(10).join(current))
            current = [line]
        else:
            current.append(line)
    if current:
        chunks.append(chr(10).join(current))
    return chunks

def chunk_email(text, max_chars=1000):
    lines = text.split(chr(10))
    kept = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith(">"):
            break
        if re.match(r"^On .+ wrote:\s*$", stripped):
            break
        kept.append(line)
    body = chr(10).join(kept).strip()
    return [body] if body else []

NATURE_CATEGORIES = ["prose", "structured", "faq", "legal", "transcript", "procedural", "sectioned", "email"]

CLASSIFY_SYSTEM_PROMPT = """You classify a piece of text by its structural nature, not its file format.
Respond with EXACTLY ONE WORD from this list, nothing else:

prose - free-flowing paragraphs, narrative or descriptive writing, no strong internal structure (e.g. articles, essays, general documents)
structured - dense factual/tabular content made of legible, meaningful words or numbers in short key:value-like lines (e.g. CSV rows, config settings, spec tables), no narrative flow
faq - a series of distinct question-and-answer pairs
legal - numbered sections, articles, or clauses (contracts, policies, terms)
transcript - a back-and-forth dialogue with named speaker turns
procedural - a numbered sequence of steps or instructions to follow in order
sectioned - short headers each followed by related content (manuals, resumes, spec sheets)
email - message content with quoted reply chains or "On [date] wrote:" boundaries
none_of_these - does not clearly fit any of the above, including garbled text, corrupted encoding, random symbols, or non-linguistic noise that is not coherent human-readable content

Important: only choose a category above none_of_these if the text is coherent, legible, human-readable content. If the text is symbolic noise, gibberish, or otherwise not real readable content, respond none_of_these even if it superficially resembles punctuation or formatting from another category (e.g. stray colons or brackets do not make noise "structured")."""

def classify_nature(text_sample, model="gpt-5-nano"):
    user_msg = f"Text to classify:\n---\n{text_sample[:2000]}\n---\n\nOne word answer:"
    try:
        r = requests.post("http://litellm:4000/v1/chat/completions", json={
            "model": model,
            "messages": [
                {"role": "system", "content": CLASSIFY_SYSTEM_PROMPT},
                {"role": "user", "content": user_msg},
            ],
            "max_tokens": 500,
        }, timeout=30)
        raw = r.json()["choices"][0]["message"]["content"].strip().lower()
    except Exception as e:
        print(f"classify_nature: call failed ({e}), defaulting to none_of_these")
        return "none_of_these"
    if not raw:
        print("classify_nature: empty content from model (reasoning budget exhausted?), defaulting to none_of_these")
        return "none_of_these"
    for cat in NATURE_CATEGORIES + ["none_of_these"]:
        if cat in raw:
            return cat
    print(f"classify_nature: unrecognized reply {raw!r}, defaulting to none_of_these")
    return "none_of_these"

def fast_classify(text):
    """
    Cheap, regex-based pre-check for obviously-shaped content, so classify_nature()
    (an AI call) only gets used when the shape genuinely isn't obvious. Deliberately
    conservative: only returns a category when confident, otherwise returns None so
    the caller falls through to the AI classifier. Does not attempt to detect prose,
    structured, or sectioned -- those boundaries are too fuzzy for a cheap heuristic.
    """
    sample = text[:3000]
    lines = [l for l in sample.split(chr(10)) if l.strip()]

    faq_hits = sum(
        1 for l in lines
        if re.match(r"^\s*(Q:|Question:|A:|Answer:)", l, re.IGNORECASE) or l.strip().endswith("?")
    )
    if faq_hits >= 3:
        return "faq"

    legal_hits = len(re.findall(r"^\s*(?:Section\s+\d+|Article\s+\d+)\b", sample, re.MULTILINE))
    if legal_hits >= 2:
        return "legal"

    procedural_hits = len(re.findall(r"^\s*(?:Step\s+\d+|\d+[.)]\s)", sample, re.MULTILINE))
    if procedural_hits >= 3:
        return "procedural"

    has_email_headers = (
        bool(re.search(r"^From:\s", sample, re.MULTILINE))
        and bool(re.search(r"^(To|Subject):\s", sample, re.MULTILINE))
    )
    has_quote_marker = bool(re.search(r"^On .+ wrote:\s*$", sample, re.MULTILINE))
    if has_email_headers or has_quote_marker:
        return "email"

    transcript_hits = len(re.findall(
        r"^(?!(?:From|To|Cc|Bcc|Subject|Date|Reply-To|Product|Price|Stock|Category|Name|Amount|Total|Status|Type|Q|A|Question|Answer|ID):\s)[A-Z][A-Za-z .]{0,30}:\s",
        sample, re.MULTILINE
    ))
    if transcript_hits >= 3:
        return "transcript"

    return None

STRATEGY_PROMPT = """The following text does not fit any standard document structure category.
Suggest a Python regular expression pattern that could be used with re.split(pattern, text, flags=re.MULTILINE)
to break it into meaningful, self-contained chunks based on whatever natural boundaries exist in this
specific text (repeated symbols, blank lines, numbering, indentation changes, etc.).

Respond with ONLY the raw regex pattern, nothing else -- no code fences, no explanation, no quotes around it.
If you cannot identify any reasonable split boundary, respond with exactly: NONE

Text:
---
{text}
---

Regex pattern:"""

def generate_chunk_pattern(text_sample, model="gpt-5-nano"):
    prompt = STRATEGY_PROMPT.format(text=text_sample[:2000])
    try:
        r = requests.post("http://litellm:4000/v1/chat/completions", json={
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 500,
        }, timeout=30)
        raw = r.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        print(f"generate_chunk_pattern: call failed ({e})")
        return None
    if not raw or raw.upper() == "NONE":
        return None
    return raw

def apply_split_pattern(text, pattern, max_chars=1000):
    try:
        compiled = re.compile(pattern, flags=re.MULTILINE)
    except re.error as e:
        print(f"apply_split_pattern: invalid regex ({e})")
        return None
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(compiled.split, text)
            parts = future.result(timeout=5)
    except concurrent.futures.TimeoutError:
        print("apply_split_pattern: regex split timed out, aborting")
        return None
    except Exception as e:
        print(f"apply_split_pattern: split failed ({e})")
        return None
    parts = [p.strip() for p in parts if p and p.strip()]
    if not parts:
        return None
    chunks, current = [], ""
    for part in parts:
        if len(current) + len(part) <= max_chars:
            current += part + chr(10)
        else:
            if current:
                chunks.append(current.strip())
            current = part + chr(10)
    if current:
        chunks.append(current.strip())
    return chunks

def sanity_check_chunks(chunks, original_text, max_chars=1000):
    if not chunks:
        return False
    if any(not c.strip() for c in chunks):
        return False
    orig_len = len(original_text.strip())
    if orig_len == 0:
        return False
    total_len = sum(len(c) for c in chunks)
    if total_len < orig_len * 0.5:
        return False
    if any(len(c) > max_chars * 3 for c in chunks):
        return False
    return True

MAX_STRATEGY_RETRIES = 2

def chunk_none_of_these(text, max_chars=1000):
    for attempt in range(1, MAX_STRATEGY_RETRIES + 1):
        pattern = generate_chunk_pattern(text)
        if pattern is None:
            print(f"chunk_none_of_these: attempt {attempt} - no pattern suggested")
            break
        chunks = apply_split_pattern(text, pattern, max_chars=max_chars)
        if chunks and sanity_check_chunks(chunks, text, max_chars=max_chars):
            print(f"chunk_none_of_these: attempt {attempt} succeeded")
            return chunks
        print(f"chunk_none_of_these: attempt {attempt} failed sanity check")
    print("chunk_none_of_these: falling back to chunk_text")
    return chunk_text(text, max_chars=max_chars)

def classify_and_chunk(text):
    nature = fast_classify(text)
    if nature:
        print(f"classify_and_chunk: fast-classified as \'{nature}\' (no AI call needed)")
    else:
        nature = classify_nature(text)
        print(f"classify_and_chunk: classified as \'{nature}\'")
    if nature in CHUNKER_DISPATCH:
        return CHUNKER_DISPATCH[nature](text)
    return chunk_none_of_these(text)


CONTEXT_SYSTEM_TEMPLATE = """You are given the full text of a document. Write a short 1-2 sentence context note (under 50 words) situating the following chunk within the document -- what section/topic it covers and what the document overall is about. Respond with ONLY the context note, no quotes, no repeating the chunk itself.

Full document:
---
{doc}
---"""

def generate_chunk_context(doc_text, chunk, model="gpt-5-nano"):
    # Contextual retrieval (Anthropic's technique): a short AI-written note is
    # prepended to each chunk before embedding, so isolated chunks (a bare
    # table row, a short paragraph) carry some of the surrounding document's
    # meaning instead of standing completely alone. The document goes in the
    # system message and only the chunk varies in the user message across
    # repeated calls for the same document, so provider-side prompt caching
    # can reduce the cost of contextualizing many chunks from one source.
    system_msg = CONTEXT_SYSTEM_TEMPLATE.format(doc=doc_text[:6000])
    user_msg = f"Chunk to contextualize:\n---\n{chunk[:1500]}\n---\n\nContext note:"
    try:
        r = requests.post("http://litellm:4000/v1/chat/completions", json={
            "model": model,
            "messages": [
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg},
            ],
            "max_tokens": 1000,
        }, timeout=30)
        raw = r.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        print(f"generate_chunk_context: call failed ({e}), storing chunk without context")
        return None
    if not raw:
        print("generate_chunk_context: empty content from model, storing chunk without context")
        return None
    return raw

def embed(text):
    r = requests.post('http://litellm:4000/v1/embeddings', json={'model': 'embed-small', 'input': text})
    return r.json()['data'][0]['embedding']

def store(source, chunks, doc_text=None):
    conn = psycopg2.connect(host='postgres', dbname='litellm', user='litellm', password=os.environ.get('POSTGRES_PASSWORD'))
    cur = conn.cursor()
    for i, chunk in enumerate(chunks):
        text_to_store = chunk
        if doc_text:
            context = generate_chunk_context(doc_text, chunk)
            if context:
                text_to_store = f'{context}\n\n{chunk}'
        vector = embed(text_to_store)
        cur.execute('INSERT INTO rag_chunks (source, content, embedding) VALUES (%s, %s, %s)', (source, text_to_store, vector))
        print(f'Stored chunk {i+1}/{len(chunks)} [{source}]')
    conn.commit()
    conn.close()

def ingest_wiki(topic):
    print(f'Fetching: {topic}')
    resp = requests.get(
        'https://en.wikipedia.org/w/api.php',
        headers={'User-Agent': 'HermesRAGTest/1.0 (research project)'},
        params={'action': 'query', 'format': 'json', 'titles': topic, 'prop': 'extracts', 'explaintext': True}
    )
    pages = resp.json()['query']['pages']
    page = list(pages.values())[0]
    if 'missing' in page:
        print(f'No Wikipedia article found for "{topic}"')
        sys.exit(1)
    raw_text = page.get('extract', '')
    print(f'Fetched {len(raw_text)} characters')
    chunks = classify_and_chunk(raw_text)
    print(f'Created {len(chunks)} chunks')
    store(f'wikipedia:{topic}', chunks, doc_text=raw_text)

def ingest_pdf(path):
    # pdfplumber handles both text and table extraction. Tables are detected
    # and pulled out separately (via chunk_structured, same treatment as docx
    # tables) so their rows/columns don't get jumbled into the flat page text --
    # pypdf had no concept of table structure at all, this was the "Still Open"
    # gap noted in the summary doc.
    import pdfplumber
    fname = os.path.basename(path)
    total = 0
    ocr_pages = 0
    table_total = 0
    tables_found = 0
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            try:
                tables = page.find_tables()
            except Exception as e:
                print(f'  Page {page_num}: table detection error: {e}')
                tables = []
            table_bboxes = [t.bbox for t in tables]

            try:
                if table_bboxes:
                    def not_within_bboxes(obj, bboxes=table_bboxes):
                        v_mid = (obj["top"] + obj["bottom"]) / 2
                        h_mid = (obj["x0"] + obj["x1"]) / 2
                        return not any(
                            (h_mid >= x0) and (h_mid < x1) and (v_mid >= top) and (v_mid < bottom)
                            for x0, top, x1, bottom in bboxes
                        )
                    text = page.filter(not_within_bboxes).extract_text() or ''
                else:
                    text = page.extract_text() or ''
            except Exception as e:
                print(f'  Page {page_num}: text extraction error: {e}')
                text = ''

            if not text.strip():
                # No text layer -- likely a scanned page. Fall back to rendering
                # the page as an image and running it through Tesseract OCR.
                try:
                    from pdf2image import convert_from_path
                    import pytesseract
                    images = convert_from_path(path, first_page=page_num, last_page=page_num)
                    if images:
                        text = pytesseract.image_to_string(images[0])
                        if text.strip():
                            ocr_pages += 1
                            print(f'  Page {page_num}: no text layer, OCR fallback succeeded')
                except Exception as e:
                    print(f'  Page {page_num}: OCR fallback failed: {e}')
                    text = ''

            if text.strip():
                chunks = classify_and_chunk(text)
                if chunks:
                    store(f'{fname}#page={page_num}', chunks, doc_text=text)
                    total += len(chunks)
            elif not table_bboxes:
                print(f'  Page {page_num}: no extractable text even after OCR, skipping')

            for t_idx, table in enumerate(tables, start=1):
                try:
                    rows = table.extract()
                except Exception as e:
                    print(f'  Page {page_num}, table {t_idx}: extraction error: {e}')
                    continue
                if not rows:
                    continue
                row_lines = [
                    ' | '.join((cell or '').strip() for cell in row)
                    for row in rows if any(cell and cell.strip() for cell in row)
                ]
                table_text = chr(10).join(row_lines)
                if table_text.strip():
                    t_chunks = chunk_structured(table_text)
                    if t_chunks:
                        store(f'{fname}#page={page_num}#table={t_idx}', t_chunks, doc_text=(text + chr(10) + table_text) if text.strip() else table_text)
                        table_total += len(t_chunks)
                        tables_found += 1
    print(f'Done. {total} chunks stored from {fname} ({ocr_pages} pages used OCR fallback, '
          f'{tables_found} tables found producing {table_total} additional chunks)')

def ingest_docx(path):
    import docx
    d = docx.Document(path)
    fname = os.path.basename(path)
    full_text = chr(10).join(p.text for p in d.paragraphs)
    chunks = classify_and_chunk(full_text)
    if chunks:
        store(fname, chunks, doc_text=full_text)
    table_rows = []
    for table in d.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                table_rows.append(row_text)
    if table_rows:
        t_chunks = chunk_structured(chr(10).join(table_rows))
        if t_chunks:
            store(f'{fname}#tables', t_chunks, doc_text=full_text)
    print(f'Done processing {fname}')

def ingest_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    fname = os.path.basename(path)
    total = 0
    slide_texts = []
    for slide in prs.slides:
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        slide_texts.append(chr(10).join(t for t in texts if t.strip()))
    full_deck_text = chr(10).join(t for t in slide_texts if t.strip())
    for slide_num, slide_text in enumerate(slide_texts, start=1):
        if not slide_text.strip():
            continue
        chunks = classify_and_chunk(slide_text)
        if chunks:
            store(f'{fname}#slide={slide_num}', chunks, doc_text=full_deck_text)
            total += len(chunks)
    print(f'Done. {total} chunks stored from {fname}')

def ingest_txt(path):
    fname = os.path.basename(path)
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        raw_text = f.read()
    chunks = classify_and_chunk(raw_text)
    store(fname, chunks, doc_text=raw_text)
    print(f'Done. {len(chunks)} chunks stored from {fname}')

def ingest_image(path):
    import pytesseract
    from PIL import Image
    fname = os.path.basename(path)
    try:
        img = Image.open(path)
        full_text = pytesseract.image_to_string(img)
    except Exception as e:
        print(f'Failed to OCR {fname}: {e}')
        return
    if not full_text.strip():
        print(f'No text extracted via OCR from {fname}')
        return
    chunks = chunk_structured(full_text)
    if chunks:
        store(fname, chunks, doc_text=full_text)
    print(f'Done. {len(chunks)} chunks stored from {fname}')

def ingest_file(path):
    if not os.path.isfile(path):
        print(f'File not found: {path}')
        sys.exit(1)
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        ingest_pdf(path)
    elif ext == '.docx':
        ingest_docx(path)
    elif ext == '.pptx':
        ingest_pptx(path)
    elif ext == '.txt':
        ingest_txt(path)
    elif ext in ('.png', '.jpg', '.jpeg'):
        ingest_image(path)
    else:
        print(f'Unsupported file type: {ext} (supported: .pdf, .docx, .pptx, .txt, .png, .jpg, .jpeg)')
        sys.exit(1)

def ingest_db(query, column, label):
    conn = psycopg2.connect(host='postgres', dbname='litellm', user='litellm', password=os.environ.get('POSTGRES_PASSWORD'))
    cur = conn.cursor()
    cur.execute(query)
    col_names = [desc[0] for desc in cur.description]
    if column not in col_names:
        print(f'Column "{column}" not in result columns: {col_names}')
        sys.exit(1)
    col_idx = col_names.index(column)
    rows = cur.fetchall()
    conn.close()
    print(f'Query returned {len(rows)} rows')
    total = 0
    for row_num, row in enumerate(rows, start=1):
        text = row[col_idx]
        if not text or not str(text).strip():
            continue
        chunks = classify_and_chunk(str(text))
        if chunks:
            row_context = f"This is row {row_num} from a database table labeled '{label}'.\n\n{text}"
            store(f'{label}#row={row_num}', chunks, doc_text=row_context)
            total += len(chunks)
    print(f'Done. {total} chunks stored from db:{label}')

def chunk_structured(text, max_chars=800):
    # For structured key:value rows (dbrow). Unlike chunk_text, this keeps
    # every line regardless of length -- short fields like price/category
    # are legitimate content here, not prose noise to filter out.
    lines = [l.strip() for l in text.split(chr(10)) if l.strip()]
    chunks, current = [], ''
    for line in lines:
        if len(current) + len(line) <= max_chars:
            current += line + chr(10)
        else:
            if current:
                chunks.append(current.strip())
            current = line + chr(10)
    if current:
        chunks.append(current.strip())
    return chunks

CHUNKER_DISPATCH = {
    "prose": chunk_text,
    "structured": chunk_structured,
    "faq": chunk_faq,
    "legal": chunk_legal,
    "transcript": chunk_transcript,
    "procedural": chunk_procedural,
    "sectioned": chunk_sectioned,
    "email": chunk_email,
}


def ingest_db_multi(query, label):
    # Handles tables where relevant info is spread across multiple columns
    # (e.g. products: name, price, category, description) rather than one
    # dedicated text column. Every returned column is included, labeled, so
    # the embedding gets full row context instead of a bare value.
    conn = psycopg2.connect(host='postgres', dbname='litellm', user='litellm', password=os.environ.get('POSTGRES_PASSWORD'))
    cur = conn.cursor()
    cur.execute(query)
    col_names = [desc[0] for desc in cur.description]
    rows = cur.fetchall()
    conn.close()
    print(f'Query returned {len(rows)} rows, columns: {col_names}')
    total = 0
    for row_num, row in enumerate(rows, start=1):
        lines = []
        for col_name, value in zip(col_names, row):
            if value is not None and str(value).strip():
                lines.append(f'{col_name}: {value}')
        if not lines:
            continue
        row_text = chr(10).join(lines)
        chunks = chunk_structured(row_text)
        if chunks:
            row_context = f"This is row {row_num} from a database table labeled '{label}', with columns: {', '.join(col_names)}.\n\n{row_text}"
            store(f'{label}#row={row_num}', chunks, doc_text=row_context)
            total += len(chunks)
    print(f'Done. {total} chunks stored from db:{label}')

def usage():
    print('Usage:')
    print('  python3 rag_ingest.py wiki "Article Title"')
    print('  python3 rag_ingest.py file /path/to/document.pdf|.docx|.pptx|.txt')
    print('  python3 rag_ingest.py db "SELECT ... FROM ..." text_column_name label')
    print('  python3 rag_ingest.py dbrow "SELECT col1, col2, ... FROM ..." label')
    sys.exit(1)

if __name__ == '__main__':
    if len(sys.argv) < 2:
        usage()

    mode = sys.argv[1]
    if mode == 'wiki':
        if len(sys.argv) < 3:
            usage()
        ingest_wiki(sys.argv[2])
    elif mode == 'file':
        if len(sys.argv) < 3:
            usage()
        ingest_file(sys.argv[2])
    elif mode == 'db':
        if len(sys.argv) < 5:
            usage()
        ingest_db(sys.argv[2], sys.argv[3], sys.argv[4])
    elif mode == 'dbrow':
        if len(sys.argv) < 4:
            usage()
        ingest_db_multi(sys.argv[2], sys.argv[3])
    else:
        usage()
